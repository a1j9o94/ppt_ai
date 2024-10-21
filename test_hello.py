import pytest
from hello import app
import json
import os
from pptx import Presentation

@pytest.fixture
def client():
    app.config['TESTING'] = True
    with app.test_client() as client:
        yield client

def test_home_route(client):
    response = client.get('/')
    assert response.status_code == 200
    assert b"Welcome to the PowerPoint Manipulation App!" in response.data

def test_create_ppt_route(client):
    # Test with valid name and title
    response = client.post('/create_ppt', 
                           data=json.dumps({'name': 'TestPresentation', 'title': 'Custom Title'}),
                           content_type='application/json')
    assert response.status_code == 200
    data = json.loads(response.data)
    assert "PowerPoint 'TestPresentation.pptx' created successfully with title 'Custom Title'" in data['message']
    assert os.path.exists(data['file_path'])
    
    # Verify the title in the created PowerPoint
    prs = Presentation(data['file_path'])
    assert prs.slides[0].shapes.title.text == 'Custom Title'
    
    # Clean up: remove the created file
    os.remove(data['file_path'])

    # Test with valid name but no title (should use default)
    response = client.post('/create_ppt', 
                           data=json.dumps({'name': 'DefaultTitleTest'}),
                           content_type='application/json')
    assert response.status_code == 200
    data = json.loads(response.data)
    assert "PowerPoint 'DefaultTitleTest.pptx' created successfully with title 'New Presentation'" in data['message']
    assert os.path.exists(data['file_path'])
    
    # Verify the default title in the created PowerPoint
    prs = Presentation(data['file_path'])
    assert prs.slides[0].shapes.title.text == 'New Presentation'
    
    # Clean up: remove the created file
    os.remove(data['file_path'])

    # Test with missing name
    response = client.post('/create_ppt', 
                           data=json.dumps({'title': 'Title without name'}),
                           content_type='application/json')
    assert response.status_code == 400
    data = json.loads(response.data)
    assert "Name is required" in data['error']

def test_create_ppt_file_content(client):
    response = client.post('/create_ppt', 
                           data=json.dumps({'name': 'ContentTest'}),
                           content_type='application/json')
    assert response.status_code == 200
    data = json.loads(response.data)
    
    # Check if file exists and has content (size > 0)
    assert os.path.exists(data['file_path'])
    assert os.path.getsize(data['file_path']) > 0
    
    # Clean up: remove the created file
    os.remove(data['file_path'])

def test_delete_ppt_route(client):
    # First, create a PowerPoint file
    create_response = client.post('/create_ppt', 
                                  data=json.dumps({'name': 'DeleteTest'}),
                                  content_type='application/json')
    assert create_response.status_code == 200
    create_data = json.loads(create_response.data)
    file_path = create_data['file_path']
    
    # Test deleting the file
    delete_response = client.delete('/delete_ppt', 
                                    data=json.dumps({'name': 'DeleteTest'}),
                                    content_type='application/json')
    assert delete_response.status_code == 200
    delete_data = json.loads(delete_response.data)
    assert "PowerPoint 'DeleteTest.pptx' deleted successfully" in delete_data['message']
    assert not os.path.exists(file_path)

    # Test deleting a non-existent file
    delete_response = client.delete('/delete_ppt', 
                                    data=json.dumps({'name': 'NonExistentFile'}),
                                    content_type='application/json')
    assert delete_response.status_code == 404
    delete_data = json.loads(delete_response.data)
    assert "File 'NonExistentFile.pptx' not found" in delete_data['error']

    # Test with missing name
    delete_response = client.delete('/delete_ppt', 
                                    data=json.dumps({}),
                                    content_type='application/json')
    assert delete_response.status_code == 400
    delete_data = json.loads(delete_response.data)
    assert "Name is required" in delete_data['error']
