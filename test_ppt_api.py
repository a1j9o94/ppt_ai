import pytest
import os
from pptx import Presentation as PPTXPresentation
from ppt_helpers import create_ppt, delete_ppt, add_slide, create_ppt_from_json
from models import PresentationModel, SlideModel, Section, LayoutType, PresentationRequest
from unittest.mock import Mock
from ppt_creation_agent import (
    create_presentation_from_prompt,
    create_presentation,
    add_slide,
    save_presentation,
    delete_presentation
)
import logging

def test_create_ppt():
    # Test with valid name and title
    file_path = create_ppt('TestPresentation', 'Custom Title')
    assert os.path.exists(file_path)
    
    # Clean up: remove the created file
    os.remove(file_path)

    # Test with missing name
    with pytest.raises(ValueError, match="Name is required"):
        create_ppt('')

def test_create_ppt_file_content():
    file_path = create_ppt('ContentTest')
    
    # Check if file exists and has content (size > 0)
    assert os.path.exists(file_path)
    assert os.path.getsize(file_path) > 0
    
    # Clean up: remove the created file
    os.remove(file_path)

def test_create_ppt_with_slides():
    file_path = create_ppt('SlidesTest', slides=[{'slide_title': 'Slide 1'}, {'slide_title': 'Slide 2'}])
    assert os.path.exists(file_path)
    os.remove(file_path)

def test_delete_ppt():
    # First, create a PowerPoint file
    file_path = create_ppt('DeleteTest')
    assert os.path.exists(file_path)
    
    # Test deleting the file
    delete_ppt('DeleteTest')
    assert not os.path.exists(file_path)

    # Test deleting a non-existent file
    with pytest.raises(FileNotFoundError, match="File 'NonExistentFile.pptx' not found"):
        delete_ppt('NonExistentFile')

    # Test with missing name
    with pytest.raises(ValueError, match="Name is required"):
        delete_ppt('')

def test_add_slide():
    # First, create a PowerPoint file
    file_path = create_ppt('AddSlideTest', 'Initial Title')
    assert os.path.exists(file_path)

    # Test adding a new slide
    add_slide('AddSlideTest', 'New Slide')

    # Verify the new slide in the PowerPoint
    prs = PPTXPresentation(file_path)
    assert len(prs.slides) == 2  # Initial slide + new slide
    assert prs.slides[1].shapes.title.text == 'New Slide'

    # Test adding a slide to non-existent file
    error = add_slide('NonExistentFile', 'New Slide')
    print(error)
    assert "error" in error
    assert "not found" in error["error"]

    # Test with missing name
    error = add_slide('', 'New Slide')
    assert "error" in error
    assert "ppt_name is required" in error["error"]

    # Test with missing slide_title
    error = add_slide('AddSlideTest', '')
    assert "error" in error
    assert "slide_title is required" in error["error"]

    # Clean up: remove the created file
    os.remove(file_path)

def test_create_ppt_from_json():
    # Test JSON data using Pydantic models
    test_presentation = PresentationRequest(
        presentation=PresentationModel(
            name="Test_Presentation",
            title="Test Title",
            slides=[
                SlideModel(
                    slide_title="First Slide",
                    ppt_name="Test_Presentation",
                    sections=[
                        Section(header="Section 1", content=["Point 1", "Point 2"]),
                        Section(header="Section 2", content=["Point A", "Point B"])
                    ]
                ),
                SlideModel(
                    slide_title="Second Slide",
                    ppt_name="Test_Presentation",
                    layout=LayoutType.COLUMN,
                    columns=2,
                    sections=[
                        Section(header="Left Column", content=["Content 1"]),
                        Section(header="Right Column", content=["Content 2"])
                    ]
                )
            ]
        )
    )

    file_path = None
    try:
        # Create presentation from JSON
        file_path = create_ppt_from_json(test_presentation.model_dump())
        
        # Verify the file exists
        assert os.path.exists(file_path)
        
        # Verify presentation content
        prs = PPTXPresentation(file_path)
        assert len(prs.slides) == 3  # Title slide + 2 content slides
        
        # Verify slide titles
        assert prs.slides[1].shapes.title.text == "First Slide"
        assert prs.slides[2].shapes.title.text == "Second Slide"

    finally:
        # Clean up
        if file_path and os.path.exists(file_path):
            os.remove(file_path)

def test_create_ppt_from_complex_json():
    # Test with the full IBM Pricing Comeback example using Pydantic models
    complex_presentation = PresentationRequest(
        presentation=PresentationModel(
            name="IBM_Test_Pricing",
            title="IBM Pricing Test - 2024",
            slides=[
                SlideModel(
                    slide_title="Meeting Objectives",
                    sections=[
                        Section(header="Objective 1", content=["Test objective 1"]),
                        Section(header="Objective 2", content=["Test objective 2"])
                    ]
                ),
                SlideModel(
                    slide_title="Pricing Overview",
                    layout=LayoutType.COLUMN,
                    columns=2,
                    sections=[
                        Section(header="Left Side", content=["Test content left"]),
                        Section(header="Right Side", content=["Test content right"])
                    ]
                )
            ]
        )
    )

    file_path = None
    try:
        file_path = create_ppt_from_json(complex_presentation.model_dump())
        
        # Verify the file exists
        assert os.path.exists(file_path)
        
        # Verify presentation content
        prs = PPTXPresentation(file_path)
        assert len(prs.slides) == 3  # Title slide + 2 content slides
        
        # Verify specific slide titles
        assert prs.slides[1].shapes.title.text == "Meeting Objectives"
        assert prs.slides[2].shapes.title.text == "Pricing Overview"

    finally:
        # Clean up
        if file_path and os.path.exists(file_path):
            os.remove(file_path)

@pytest.fixture
def mock_openai_response():
    """Mock OpenAI API response"""
    class MockMessage:
        def __init__(self, function_call=None):
            self.function_call = function_call
            self.content = "Test response"

    class MockChoice:
        def __init__(self, message):
            self.message = message

    class MockResponse:
        def __init__(self, choices):
            self.choices = choices

    return MockResponse([MockChoice(MockMessage())])

def test_create_presentation():
    # Test successful creation
    result = create_presentation("test_pres", "Test Presentation")
    assert "file_path" in result
    assert "title" in result
    
    # Clean up
    if "file_path" in result:
        os.remove(result["file_path"])

    # Test error handling
    result = create_presentation("", "")
    assert "error" in result

def test_add_slide_function():
    # First create a presentation
    create_presentation("test_slides", "Test Slides")
    
    # Test adding a simple slide
    result = add_slide(ppt_name="test_slides", slide_title="Simple Slide")
    assert "message" in result
    assert "error" not in result
    
    # Test adding a column layout slide
    sections = [
        {"header": "Col 1", "content": ["Point 1"]},
        {"header": "Col 2", "content": ["Point 2"]}
    ]
    result = add_slide(
        ppt_name="test_slides",
        slide_title="Column Slide",
        layout="COLUMN",
        sections=sections
    )
    assert "message" in result
    assert "error" not in result
    
    # Test adding a row layout slide
    result = add_slide(
        ppt_name="test_slides",
        slide_title="Row Slide",
        layout="ROW",
        sections=sections
    )
    assert "message" in result
    assert "error" not in result
    
    # Test error case
    result = add_slide(ppt_name="nonexistent", slide_title="Error Slide")
    assert "error" in result
    assert "not found" in result["error"]
    
    # Test with missing ppt_name
    result = add_slide(ppt_name="", slide_title="New Slide")
    assert "error" in result
    assert "ppt_name is required" in result["error"]

    # Test with missing slide_title
    result = add_slide(ppt_name="test_slides", slide_title="")
    assert "error" in result
    assert "slide_title is required" in result["error"]
    
    # Clean up
    delete_presentation("test_slides")

def test_save_presentation():
    # Create a presentation first
    create_presentation("test_save", "Test Save")
    
    # Test saving
    result = save_presentation("test_save")
    assert "file_path" in result
    assert "error" not in result
    
    # Test error case
    result = save_presentation("nonexistent")
    assert "error" in result
    assert "does not exist" in result["error"]
    
    # Clean up
    delete_presentation("test_save")

def test_delete_presentation_function():
    # Create a presentation first
    create_presentation("test_delete", "Test Delete")
    
    # Test successful deletion
    result = delete_presentation("test_delete")
    assert "message" in result
    assert "error" not in result
    
    # Test deleting non-existent presentation
    result = delete_presentation("nonexistent")
    assert "error" in result

def create_mock_response(function_call_name=None, function_call_args=None, message_content=None, finish_reason=None):
    function_call = None
    if function_call_name:
        function_call = Mock(
            name=function_call_name,
            arguments=function_call_args
        )
    message = Mock(
        role='assistant',
        content=message_content,
        function_call=function_call
    )
    choice = Mock(
        message=message,
        finish_reason=finish_reason
    )
    response = Mock(
        choices=[choice]
    )
    return response

def test_create_presentation_from_prompt_full_flow():
    # Create mock responses
    responses = [
        # First response: create_presentation
        create_mock_response(
            function_call_name="create_presentation",
            function_call_args='{"name": "test_pres", "title": "Test Presentation"}',
            finish_reason="function_call"
        ),
        # Second response: add_slide
        create_mock_response(
            function_call_name="add_slide",
            function_call_args='{"ppt_name": "test_pres", "slide_title": "Introduction", "sections": [{"header": "Overview", "content": ["Point 1"]}]}',
            finish_reason="function_call"
        ),
        # Third response: save_presentation
        create_mock_response(
            function_call_name="save_presentation",
            function_call_args='{"name": "test_pres"}',
            finish_reason="function_call"
        ),
        # Fourth response: final assistant message
        create_mock_response(
            message_content="Presentation created successfully",
            finish_reason="stop"
        )
    ]

    # Create a mock client
    mock_client = Mock()
    mock_client.chat.completions.create.side_effect = responses

    # Run the function with the mock client
    prompt = "Create a presentation on AI"
    result = create_presentation_from_prompt(prompt, client=mock_client)

    # Assertions
    assert result["status"] == "completed"
    assert "file_path" in result

    # Clean up if necessary
    if "file_path" in result and os.path.exists(result["file_path"]):
        os.remove(result["file_path"])

def test_create_presentation_from_prompt():
    # Create mock responses
    responses = [
        # First response: create_presentation
        create_mock_response(
            function_call_name="create_presentation",
            function_call_args='{"name": "test_pres", "title": "Test Presentation"}',
            finish_reason="function_call"
        ),
        # Final response: assistant message
        create_mock_response(
            message_content="Presentation created successfully",
            finish_reason="stop"
        )
    ]

    # Create a mock client
    mock_client = Mock()
    mock_client.chat.completions.create.side_effect = responses

    # Run the function with the mock client
    prompt = "Create a presentation about AI"
    result = create_presentation_from_prompt(prompt, client=mock_client)

    # Assertions
    assert result["status"] == "completed"
    assert "file_path" in result

    # Clean up if necessary
    if "file_path" in result and os.path.exists(result["file_path"]):
        os.remove(result["file_path"])