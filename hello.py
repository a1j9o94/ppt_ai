from flask import Flask, render_template, request, jsonify
from pptx import Presentation
from pptx.util import Inches, Pt
import os

app = Flask(__name__)

@app.route('/')
def home():
    return "Welcome to the PowerPoint Manipulation App!"

@app.route('/create_ppt', methods=['POST'])
def create_ppt():
    name = request.json.get('name')
    title = request.json.get('title', 'New Presentation')  # Default title if not provided
    if not name:
        return jsonify({"error": "Name is required"}), 400

    # Create an empty presentation
    prs = Presentation('bain_template.pptx')

    # Add a title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title_placeholder = slide.shapes.title
    subtitle_placeholder = slide.placeholders[1]

    title_placeholder.text = title
    subtitle_placeholder.text = "Created with PowerPoint Manipulation App"

    # Save the presentation
    file_name = f"{name}.pptx"
    file_path = os.path.join(os.getcwd(), file_name)
    prs.save(file_path)

    return jsonify({
        "message": f"PowerPoint '{file_name}' created successfully with title '{title}'",
        "file_path": file_path
    })

@app.route('/delete_ppt', methods=['DELETE'])
def delete_ppt():
    name = request.json.get('name')
    if not name:
        return jsonify({"error": "Name is required"}), 400

    file_name = f"{name}.pptx"
    file_path = os.path.join(os.getcwd(), file_name)

    if not os.path.exists(file_path):
        return jsonify({"error": f"File '{file_name}' not found"}), 404

    try:
        os.remove(file_path)
        return jsonify({"message": f"PowerPoint '{file_name}' deleted successfully"})
    except Exception as e:
        return jsonify({"error": f"Failed to delete '{file_name}': {str(e)}"}), 500

if __name__ == "__main__":
    app.run(debug=True)
