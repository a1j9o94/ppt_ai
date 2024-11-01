from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from models import LayoutType
import os
import logging

# Set up logging with more detailed format
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

def create_ppt(name, title='New Presentation', slides=None):
    logger.debug("Starting create_ppt function")
    try:
        logger.debug(f"Received parameters - name: {name}, title: {title}, slides: {slides}")

        if not name:
            logger.warning("Missing required name parameter")
            raise ValueError("Name is required")

        logger.debug("Creating new presentation from template")
        try:
            prs = Presentation('ppt_templates/bain_template.pptx')
        except Exception as e:
            logger.error(f"Failed to load template: {str(e)}")
            raise

        logger.debug("Adding title slide")
        try:
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title_placeholder = slide.shapes.title
            subtitle_placeholder = slide.placeholders[1]

            title_placeholder.text = title
            subtitle_placeholder.text = "Created with PowerPoint Manipulation App"
        except Exception as e:
            logger.error(f"Failed to create title slide: {str(e)}")
            raise

        # Save initial presentation with title slide
        file_name = f"{name}.pptx"
        file_path = os.path.join(os.getcwd(), "output", file_name)
        logger.debug(f"Saving initial presentation to: {file_path}")
        try:
            prs.save(file_path)
        except Exception as e:
            logger.error(f"Failed to save initial presentation: {str(e)}")
            raise

        # Add additional slides if provided
        if slides:
            logger.debug(f"Adding {len(slides)} additional slides")
            for slide_data in slides:
                try:
                    slide_title = slide_data.get('slide_title')
                    layout = slide_data.get('layout')
                    columns = slide_data.get('columns')
                    rows = slide_data.get('rows') 
                    sections = slide_data.get('sections', [])

                    if layout == LayoutType.COLUMN.value and columns:
                        add_slide(name, slide_title, columns=columns, sections=sections)
                    elif layout == LayoutType.ROW.value and rows:
                        add_slide(name, slide_title, rows=rows, sections=sections)
                    else:
                        if sections:
                            add_slide(name, slide_title, columns=len(sections), sections=sections)
                        else:
                            add_slide(name, slide_title)
                except Exception as e:
                    logger.error(f"Failed to add slide {slide_title}: {str(e)}")
                    raise

        logger.debug("Presentation created successfully")
        return file_path
    except Exception as e:
        logger.error(f"Unexpected error in create_ppt: {str(e)}", exc_info=True)
        raise

def delete_ppt(name):
    logger.debug("Starting delete_ppt function")
    try:
        logger.debug(f"Received parameter - name: {name}")

        if not name:
            logger.warning("Missing required name parameter")
            raise ValueError("Name is required")

        file_name = f"{name}.pptx"
        file_path = os.path.join(os.getcwd(), "output", file_name)
        logger.debug(f"Attempting to delete file: {file_path}")

        if not os.path.exists(file_path):
            logger.warning(f"File not found: {file_path}")
            raise FileNotFoundError(f"File '{file_name}' not found")

        try:
            os.remove(file_path)
            logger.debug(f"Successfully deleted file: {file_path}")
        except Exception as e:
            logger.error(f"Failed to delete file: {str(e)}")
            raise
    except Exception as e:
        logger.error(f"Unexpected error in delete_ppt: {str(e)}", exc_info=True)
        raise

def add_slide(name, slide_title, columns=None, rows=None, sections=[]):
    logger.debug("Starting add_slide function")
    try:
        logger.debug(f"Received parameters - name: {name}, slide_title: {slide_title}")

        if not name:
            logger.warning("Missing required name parameter")
            raise ValueError("ppt_name is required")
        
        if not slide_title:
            logger.warning("Missing required slide_title parameter")
            raise ValueError("slide_title is required")

        logger.debug(f"Layout parameters - columns: {columns}, rows: {rows}, sections: {len(sections)}")

        file_name = f"{name}.pptx"
        file_path = os.path.join(os.getcwd(), "output", file_name)
        logger.debug(f"File path: {file_path}")

        if not os.path.exists(file_path):
            logger.warning(f"File not found: {file_path}")
            raise FileNotFoundError(f"File '{file_name}' not found")

        # Validate layout parameters
        if columns and rows:
            logger.warning("Both rows and columns specified")
            raise ValueError("Cannot specify both rows and columns")
        
        if columns and len(sections) != columns:
            logger.warning(f"Section count mismatch - expected {columns}, got {len(sections)}")
            raise ValueError(f"Number of sections must match number of columns ({columns})")
        
        if rows and len(sections) != rows:
            logger.warning(f"Section count mismatch - expected {rows}, got {len(sections)}")
            raise ValueError(f"Number of sections must match number of rows ({rows})")

        # Convert Pydantic models to dictionaries if needed
        processed_sections = []
        for section in sections:
            if hasattr(section, 'model_dump'):
                processed_sections.append(section.model_dump())
            else:
                processed_sections.append(section)

        has_sizes = any('size' in section for section in processed_sections)
        if has_sizes:
            logger.debug("Validating section sizes")
            if not all('size' in section for section in processed_sections):
                logger.warning("Mixed sized and unsized sections")
                raise ValueError("Cannot mix sized and unsized sections")
            
            sizes = [section['size'] for section in processed_sections if section['size'] is not None]
            if sizes:
                total_size = sum(sizes)
                logger.debug(f"Total size: {total_size}")
                if total_size < 98 or total_size > 102:
                    logger.warning(f"Invalid total size: {total_size}")
                    raise ValueError("Section sizes must sum to approximately 100%")

        logger.debug("Opening presentation")
        prs = Presentation(file_path)
        logger.debug("Getting slide layout")
        slide_layout = prs.slide_layouts[1]  
        logger.debug("Adding slide")
        slide = prs.slides.add_slide(slide_layout)

        logger.debug("Adding title to slide")
        title = slide.shapes.title
        if title is None:
            logger.warning("Title placeholder not found, creating a new one")
            slide.shapes.title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
            title_frame = slide.shapes.title.text_frame
            title = slide.shapes.title
        else:
            title_frame = title.text_frame

        title.text = slide_title

        if columns:
            logger.debug(f"Creating column layout with {columns} columns")
            create_column_layout(slide, sections)
        elif rows:
            logger.debug(f"Creating row layout with {rows} rows")
            create_row_layout(slide, sections)
        else:
            logger.debug("No layout specified, adding basic slide")

        logger.debug("Saving presentation")
        prs.save(file_path)
        logger.debug("Presentation saved successfully")

    except Exception as e:
        logger.error(f"Error in add_slide: {str(e)}", exc_info=True)
        raise

def create_section_box(slide, start_x, start_y, width, height):
    """Create a gray box with border for a section"""
    logger.debug(f"Creating section box at ({start_x}, {start_y}) with size ({width}, {height})")
    try:
        shape = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(start_x),
            Inches(start_y),
            Inches(width),
            Inches(height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(240, 240, 240)
        shape.line.color.rgb = RGBColor(200, 200, 200)
        logger.debug("Section box created successfully")
        return shape
    except Exception as e:
        logger.error(f"Error creating section box: {str(e)}", exc_info=True)
        raise

def add_section_content(slide, start_x, start_y, width, height, section):
    """Add header and content to a section"""
    logger.debug(f"Adding content to section at ({start_x}, {start_y})")
    try:
        # Add header
        logger.debug(f"Adding header: {section['header']}")
        header_box = slide.shapes.add_textbox(
            Inches(start_x + 0.1),
            Inches(start_y + 0.1),
            Inches(width - 0.2),
            Inches(0.5)
        )
        header_frame = header_box.text_frame
        header_frame.text = section['header']
        header_frame.paragraphs[0].font.bold = True

        # Add content
        logger.debug(f"Adding {len(section['content'])} content items")
        content_box = slide.shapes.add_textbox(
            Inches(start_x + 0.1),
            Inches(start_y + 0.7),
            Inches(width - 0.2),
            Inches(height - 0.8)
        )
        content_frame = content_box.text_frame
        
        # Clear any existing paragraphs
        while len(content_frame.paragraphs) > 0:
            p = content_frame.paragraphs[0]
            p._p.getparent().remove(p._p)
            
        # Add content with our own bullet style
        for item in section['content']:
            p = content_frame.add_paragraph()
            p.text = f"• {item}"
            p.level = 0  # Base level for bullets
            # Remove any default bullets
            p.bullet = None
            
        logger.debug("Section content added successfully")
    except Exception as e:
        logger.error(f"Error adding section content: {str(e)}", exc_info=True)
        raise

def create_column_layout(slide, sections):
    """Create a column-based layout"""
    logger.debug(f"Creating column layout with {len(sections)} sections")
    try:
        # Constants based on actual slide dimensions
        margin_x = 0.5  # Left/right margin
        margin_y = 1.5  # Top margin (space for title)
        section_margin = 0.2  # Margin between sections
        
        # Calculate total width accounting for section margins
        total_margin_space = section_margin * (len(sections) - 1)  # Space needed for margins between sections
        total_width = 12.33 - total_margin_space  # Full width minus margins and space between sections
        section_height = 5.0  # Leave room for title and bottom margin
        
        for i, section in enumerate(sections):
            logger.debug(f"Processing column {i+1}")
            # Fix the size calculation to handle None values
            if section is None or not isinstance(section, dict):
                size_factor = 1.0 / len(sections)
            else:
                size_factor = section.get('size', 100) / 100 if section.get('size') else 1.0 / len(sections)
            
            width = total_width * size_factor
            
            # Calculate start_x using the same safe size calculation
            previous_sections_size = sum(
                (s.get('size', 100) / 100 if s is not None and isinstance(s, dict) and s.get('size') 
                 else 1.0 / len(sections)) 
                for s in sections[:i]
            )
            start_x = margin_x + (total_width * previous_sections_size) + (section_margin * i)
            
            create_section_box(slide, start_x, margin_y, int(width), int(section_height))
            add_section_content(slide, start_x, margin_y, int(width), int(section_height), section)
        logger.debug("Column layout created successfully")
    except Exception as e:
        logger.error(f"Error creating column layout: {str(e)}", exc_info=True)
        raise

def create_row_layout(slide, sections):
    """Create a row-based layout"""
    logger.debug(f"Creating row layout with {len(sections)} sections")
    try:
        # Constants based on actual slide dimensions
        margin_x = 0.5  # Left/right margin
        margin_y = 1.5  # Top margin (space for title)
        section_margin = 0.2  # Margin between sections
        
        # Calculate total height accounting for section margins
        total_margin_space = section_margin * (len(sections) - 1)
        total_width = 12.33  # Full width minus margins
        total_height = 5.0 - total_margin_space  # Available height minus space between sections
        
        for i, section in enumerate(sections):
            logger.debug(f"Processing row {i+1}")
            # Fix the size calculation to handle None values
            if section is None or not isinstance(section, dict):
                size_factor = 1.0 / len(sections)
            else:
                size_factor = section.get('size', 100) / 100 if section.get('size') else 1.0 / len(sections)
            
            height = total_height * size_factor
            
            # Calculate current_y using the same safe size calculation
            previous_sections_size = sum(
                (s.get('size', 100) / 100 if s is not None and isinstance(s, dict) and s.get('size')
                 else 1.0 / len(sections))
                for s in sections[:i]
            )
            current_y = margin_y + (total_height * previous_sections_size) + (section_margin * i)

            create_section_box(slide, margin_x, current_y, total_width, height)
            add_section_content(slide, margin_x, current_y, total_width, height, section)
            
        logger.debug("Row layout created successfully")
    except Exception as e:
        logger.error(f"Error creating row layout: {str(e)}", exc_info=True)
        raise

def create_ppt_from_json(json_data):
    """Create a complete PowerPoint presentation from JSON data"""
    logger.debug("Starting create_ppt_from_json function")
    try:
        # Handle both Pydantic models and raw dictionaries
        if hasattr(json_data, 'model_dump'):
            presentation_data = json_data.model_dump()
        else:
            presentation_data = json_data

        presentation_data = presentation_data.get('presentation')
        if not presentation_data:
            raise ValueError("JSON must contain a 'presentation' object")

        name = presentation_data.get('name')
        title = presentation_data.get('title', 'New Presentation')
        slides = presentation_data.get('slides', [])

        if not name:
            raise ValueError("Presentation name is required")

        # Create the initial presentation
        file_path = create_ppt(name, title)

        # Add each slide from the JSON
        for slide_data in slides:
            slide_title = slide_data.get('slide_title')
            if not slide_title:
                logger.warning("Skipping slide without title")
                continue

            # Check for layout specification
            layout = slide_data.get('layout')
            columns = slide_data.get('columns')
            rows = slide_data.get('rows')
            sections = slide_data.get('sections', [])

            # Add the slide with specified layout
            if layout == LayoutType.COLUMN.value and columns:
                add_slide(name, slide_title, columns=columns, sections=sections)
            elif layout == LayoutType.ROW.value and rows:
                add_slide(name, slide_title, rows=rows, sections=sections)
            else:
                # Default to columns if multiple sections present
                if sections:
                    add_slide(name, slide_title, columns=len(sections), sections=sections)
                else:
                    add_slide(name, slide_title)

        logger.info(f"Successfully created presentation from JSON: {file_path}")
        return file_path

    except Exception as e:
        logger.error(f"Error in create_ppt_from_json: {str(e)}", exc_info=True)
        raise
