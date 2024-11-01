from pptx import Presentation
from pptx.util import Inches
import logging

# Set up logging
logging.basicConfig(
    level=logging.DEBUG,
    format='%(asctime)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

def analyze_template_dimensions():
    try:
        # Load the template
        logger.info("Loading template presentation")
        prs = Presentation('bain_template.pptx')
        
        # Get the slide dimensions
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        
        # Convert to inches for easier understanding
        width_inches = slide_width / 914400  # EMU to inches conversion
        height_inches = slide_height / 914400
        
        logger.info(f"Slide dimensions:")
        logger.info(f"Width: {width_inches:.2f} inches")
        logger.info(f"Height: {height_inches:.2f} inches")
        logger.info(f"Aspect ratio: {width_inches/height_inches:.2f}:1")

        # If there are any slides in the template, analyze their content
        if len(prs.slides) > 0:
            logger.info("\nAnalyzing first slide:")
            slide = prs.slides[0]
            
            # Analyze shapes and their positions
            for shape in slide.shapes:
                if shape.has_text_frame:
                    left_inches = shape.left / 914400
                    top_inches = shape.top / 914400
                    width_inches = shape.width / 914400
                    height_inches = shape.height / 914400
                    
                    logger.info(f"\nShape with text: '{shape.text}'")
                    logger.info(f"Position: left={left_inches:.2f}\", top={top_inches:.2f}\"")
                    logger.info(f"Size: width={width_inches:.2f}\", height={height_inches:.2f}\"")

        return {
            'width': width_inches,
            'height': height_inches,
            'aspect_ratio': width_inches/height_inches
        }

    except Exception as e:
        logger.error(f"Error analyzing template: {str(e)}", exc_info=True)
        raise

if __name__ == "__main__":
    dimensions = analyze_template_dimensions()
    print(dimensions)
